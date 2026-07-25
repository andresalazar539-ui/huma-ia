# ================================================================
# huma/services/export_service.py — Exportação de relatórios
#
# Transforma o relatório de outcome (report_service.build_report) em:
#   - Planilha .xlsx  → dono manda pro contador/sócio, abre no Excel
#   - Apresentação .pptx → dono apresenta em reunião, editável
#
# Ambos respeitam as seções por meta (capabilities): só exporta o que
# existe no relatório. Identidade visual da HUMA no .pptx (papel
# quente + espresso + terracota).
# ================================================================

from datetime import datetime
from io import BytesIO

from huma.models.schemas import ClientIdentity
from huma.utils.logger import get_logger

log = get_logger("export")

# Paleta HUMA (colors_and_type.css) em RGB pro pptx/xlsx
_PAPER = "F6F2EC"
_PAPER_RAISED = "FBF8F3"
_INK = "1C1714"
_INK_3 = "6B6259"
_TERRACOTTA = "C8553D"
_SAGE_INK = "3E5540"

# Rótulos amigáveis por seção (ordem de exibição)
_SECTION_LAYOUT: list[tuple[str, str, list[tuple[str, str]]]] = [
    ("atendimento", "Atendimento", [
        ("conversas_ativas", "Conversas atendidas"),
        ("conversas_novas", "Leads novos"),
        ("fora_do_horario", "Fora do horário comercial"),
    ]),
    ("vendas", "Vendas", [
        ("receita_display", "Receita confirmada"),
        ("pagamentos", "Pagamentos aprovados"),
        ("ticket_display", "Ticket médio"),
        ("fechadas_sem_humano", "Fechadas 100% pela HUMA"),
    ]),
    ("agenda", "Agenda", [
        ("agendamentos", "Agendamentos"),
        ("realizados", "Realizados"),
        ("proximos", "Ainda por vir"),
    ]),
    ("qualificacao", "Qualificação", [
        ("leads_com_dados", "Leads com dados coletados"),
        ("enviados_crm", "Enviados pro CRM"),
        ("passados_pro_humano", "Entregues quentes pro humano"),
    ]),
    ("funil", "Funil", [
        ("descoberta", "Descobrindo"),
        ("negociando", "Negociando"),
        ("compromissados", "Compromissados"),
        ("ganhos", "Ganhos"),
        ("perdidos", "Perdidos"),
    ]),
    ("origem", "Origem dos leads", [
        ("top_conversas", "Maior origem de conversas"),
        ("top_conversoes", "Maior origem de conversões"),
    ]),
    ("follow_up", "Follow-up", [
        ("leads_reengajados", "Leads sumidos reengajados"),
        ("voltaram_a_negociar", "Voltaram a negociar"),
    ]),
]


def _periodo_label(days: int) -> str:
    return {1: "último dia", 7: "últimos 7 dias", 30: "últimos 30 dias", 90: "últimos 90 dias"}.get(
        days, f"últimos {days} dias"
    )


# ================================================================
# PLANILHA (.xlsx)
# ================================================================


def report_to_xlsx(identity: ClientIdentity, report: dict) -> bytes:
    """Planilha formatada: uma aba, seções em blocos, pronta pro Excel."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "Relatório HUMA"
    ws.column_dimensions["A"].width = 38
    ws.column_dimensions["B"].width = 22

    title_font = Font(name="Calibri", bold=True, size=16, color=_INK)
    meta_font = Font(name="Calibri", size=10, color=_INK_3)
    section_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    section_fill = PatternFill(start_color=_TERRACOTTA, end_color=_TERRACOTTA, fill_type="solid")
    label_font = Font(name="Calibri", size=11, color=_INK)
    value_font = Font(name="Calibri", bold=True, size=11, color=_INK)

    nome = identity.business_name or "Seu negócio"
    days = int(report.get("period_days", 30))

    ws["A1"] = f"Relatório HUMA — {nome}"
    ws["A1"].font = title_font
    ws["A2"] = f"Período: {_periodo_label(days)} · Gerado em {datetime.utcnow().strftime('%d/%m/%Y')}"
    ws["A2"].font = meta_font

    sections = report.get("sections", {})
    row = 4
    for key, titulo, fields in _SECTION_LAYOUT:
        data = sections.get(key)
        if not data:
            continue
        ws.cell(row=row, column=1, value=titulo).font = section_font
        ws.cell(row=row, column=1).fill = section_fill
        ws.cell(row=row, column=2).fill = section_fill
        row += 1
        for field_key, label in fields:
            if field_key not in data:
                continue
            ws.cell(row=row, column=1, value=label).font = label_font
            vcell = ws.cell(row=row, column=2, value=data[field_key])
            vcell.font = value_font
            vcell.alignment = Alignment(horizontal="right")
            row += 1
        row += 1

    # Origem: ranking completo de fontes (conversas × conversões)
    fontes = sections.get("origem", {}).get("fontes") or []
    if fontes:
        ws.cell(row=row, column=1, value="De onde vieram seus leads").font = section_font
        ws.cell(row=row, column=1).fill = section_fill
        ws.cell(row=row, column=2).fill = section_fill
        row += 1
        for f in fontes:
            ws.cell(row=row, column=1, value=f.get("origem", "")).font = label_font
            resumo = f"{f.get('conversas', 0)} conversas · {f.get('ganhos', 0)} conversões"
            if f.get("receita_cents", 0):
                resumo += f" · {f.get('receita_display', '')}"
            vcell = ws.cell(row=row, column=2, value=resumo)
            vcell.font = value_font
            vcell.alignment = Alignment(horizontal="right")
            row += 1
        row += 1

    # Inteligência: ranking de assuntos
    intel = sections.get("inteligencia", {})
    top = intel.get("top_assuntos") or []
    if top:
        ws.cell(row=row, column=1, value="O que seus leads mais pediram").font = section_font
        ws.cell(row=row, column=1).fill = section_fill
        ws.cell(row=row, column=2).fill = section_fill
        row += 1
        for i, t in enumerate(top, start=1):
            ws.cell(row=row, column=1, value=f"{i}º · {t.get('tipo', '')}").font = label_font
            vcell = ws.cell(row=row, column=2, value=t.get("vezes", 0))
            vcell.font = value_font
            vcell.alignment = Alignment(horizontal="right")
            row += 1

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ================================================================
# APRESENTAÇÃO (.pptx)
# ================================================================


def report_to_pptx(identity: ClientIdentity, report: dict) -> bytes:
    """
    Deck editável com a identidade HUMA: capa + 1 slide por seção
    (números grandes) + slide de inteligência. 16:9.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt

    paper = RGBColor.from_string(_PAPER)
    ink = RGBColor.from_string(_INK)
    ink3 = RGBColor.from_string(_INK_3)
    terracotta = RGBColor.from_string(_TERRACOTTA)
    sage = RGBColor.from_string(_SAGE_INK)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    nome = identity.business_name or "Seu negócio"
    days = int(report.get("period_days", 30))
    sections = report.get("sections", {})

    def _bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = paper

    def _text(slide, left, top, width, height, text, size, color, bold=False, italic=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        return box

    # ── Capa ──
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _text(slide, 0.9, 2.2, 11.5, 1.2, "Relatório de resultados", 40, ink, bold=True)
    _text(slide, 0.9, 3.3, 11.5, 0.9, nome, 26, terracotta, bold=True)
    _text(
        slide, 0.9, 4.2, 11.5, 0.6,
        f"{_periodo_label(days).capitalize()} · gerado pela HUMA IA em {datetime.utcnow().strftime('%d/%m/%Y')}",
        14, ink3,
    )

    # ── Slides por seção (números grandes, até 5 por slide) ──
    for key, titulo, fields in _SECTION_LAYOUT:
        data = sections.get(key)
        if not data:
            continue
        metricas = [(label, data[fk]) for fk, label in fields if fk in data]
        if not metricas:
            continue

        slide = prs.slides.add_slide(blank)
        _bg(slide)
        _text(slide, 0.9, 0.6, 11.5, 0.8, titulo, 28, ink, bold=True)

        col_w = 12.0 / min(len(metricas), 5)
        for i, (label, value) in enumerate(metricas[:5]):
            x = 0.9 + i * col_w
            destaque = key in ("vendas", "follow_up") and i in (0, 1)
            _text(slide, x, 2.4, col_w - 0.3, 1.2, str(value), 40,
                  sage if destaque else ink, bold=True)
            _text(slide, x, 3.6, col_w - 0.3, 1.0, label, 13, ink3)

    # ── Origem: ranking de fontes (conversas × conversões) ──
    fontes = (sections.get("origem", {}).get("fontes") or [])[:5]
    if fontes:
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        _text(slide, 0.9, 0.6, 11.5, 0.8, "De onde vieram seus leads", 28, ink, bold=True)
        for i, f in enumerate(fontes):
            _text(slide, 0.9, 2.0 + i * 1.0, 0.8, 0.8, f"{i + 1}º", 22, terracotta, bold=True)
            _text(slide, 1.9, 2.0 + i * 1.0, 6.5, 0.8, str(f.get("origem", "")), 20, ink)
            _text(slide, 8.6, 2.0 + i * 1.0, 2.2, 0.8, f"{f.get('conversas', 0)} conversas", 18, ink3)
            _text(slide, 10.9, 2.0 + i * 1.0, 2.0, 0.8, f"{f.get('ganhos', 0)} vendas", 18, sage)

    # ── Inteligência ──
    intel = sections.get("inteligencia", {})
    top = intel.get("top_assuntos") or []
    if top:
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        _text(slide, 0.9, 0.6, 11.5, 0.8, "O que seus leads mais pediram", 28, ink, bold=True)
        for i, t in enumerate(top):
            _text(slide, 0.9, 2.2 + i * 1.1, 0.8, 0.8, f"{i + 1}º", 24, terracotta, bold=True)
            _text(slide, 1.9, 2.2 + i * 1.1, 8.0, 0.8, str(t.get("tipo", "")), 22, ink)
            _text(slide, 10.5, 2.2 + i * 1.1, 1.8, 0.8, f"{t.get('vezes', 0)}x", 22, ink3)

    # ── Fechamento ──
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _text(slide, 0.9, 3.0, 11.5, 1.0, "Sua IA de vendas, trabalhando 24/7.", 30, ink, bold=True, italic=True)
    _text(slide, 0.9, 4.1, 11.5, 0.6, "HUMA IA · relatório gerado automaticamente", 13, ink3)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
