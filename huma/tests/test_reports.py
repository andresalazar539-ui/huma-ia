# ================================================================
# huma/tests/test_reports.py — Relatórios de outcome por meta
#
# Cobre:
#   - Seções condicionais às capabilities (vendas/agenda/qualificação)
#   - Números: receita real, funil, follow-up resgatado, fora de horário
#   - Texto do WhatsApp por meta
#   - Job: janela 8h BRT, frequência do dono, dedup, silêncio sem atividade
#   - Endpoint /reports (auth + shape)
# ================================================================

import asyncio
from datetime import datetime, timedelta

import pytest

from huma.models.schemas import ClientIdentity
from huma.services import report_service as rs


def _identity(caps) -> ClientIdentity:
    return ClientIdentity(
        client_id="cli_rep",
        business_name="Clínica Relatório",
        capabilities=caps,
        owner_phone="5511999998888",
    )


def _now_iso(days_ago=0, hour=None):
    dt = datetime.utcnow() - timedelta(days=days_ago)
    if hour is not None:
        dt = dt.replace(hour=hour, minute=0)
    return dt.isoformat()


def _fake_supa(monkeypatch, *, conversations=None, payments=None, classifications=None, clients=None):
    class Resp:
        def __init__(self, data):
            self.data = data

    class Q:
        def __init__(self, table):
            self._table = table

        def select(self, *a, **kw): return self
        def eq(self, *a): return self
        def neq(self, *a): return self
        def gte(self, *a): return self
        def lte(self, *a): return self
        def limit(self, n): return self

        def execute(self):
            data = {
                "conversations": conversations or [],
                "payments": payments or [],
                "message_classifications": classifications or [],
                "clients": clients or [],
            }.get(self._table, [])
            return Resp(data)

    class Supa:
        def table(self, name): return Q(name)

    monkeypatch.setattr(rs, "get_supabase", lambda: Supa())


_CONVS = [
    # ganho com pagamento, sem humano, com follow-up que resgatou
    {"phone": "1", "stage": "won", "created_at": _now_iso(2), "last_message_at": _now_iso(1, hour=14),
     "follow_up_count": 1, "handoff_status": "active", "lead_name_canonical": "Maria",
     "lead_email": "m@x.com", "lead_facts": ["nome: Maria"], "crm_deal_id": "d1",
     "crm_synced_at": _now_iso(1), "active_appointment_datetime": _now_iso(-1),
     "active_appointment_service": "Consulta"},
    # negociando, atendida de madrugada (fora do horário)
    {"phone": "2", "stage": "offer", "created_at": _now_iso(3), "last_message_at": _now_iso(0, hour=3),
     "follow_up_count": 0, "handoff_status": "active", "lead_name_canonical": "",
     "lead_email": "", "lead_facts": [], "crm_deal_id": "", "crm_synced_at": None,
     "active_appointment_datetime": "", "active_appointment_service": ""},
    # perdido com handoff
    {"phone": "3", "stage": "lost", "created_at": _now_iso(5), "last_message_at": _now_iso(2, hour=15),
     "follow_up_count": 2, "handoff_status": "handed_off", "lead_name_canonical": "João",
     "lead_email": "", "lead_facts": [], "crm_deal_id": "", "crm_synced_at": None,
     "active_appointment_datetime": "", "active_appointment_service": ""},
]

_PAYMENTS = [
    {"amount_cents": 35000, "paid_at": _now_iso(1), "status": "approved"},
    {"amount_cents": 15000, "paid_at": _now_iso(2), "status": "approved"},
]

_CLASSIFICATIONS = (
    [{"msg_type": "price_query"}] * 5
    + [{"msg_type": "schedule_intent"}] * 3
    + [{"msg_type": "objection"}] * 2
)


class TestBuildReportPorMeta:

    def test_cliente_que_vende_ve_vendas(self, monkeypatch):
        _fake_supa(monkeypatch, conversations=_CONVS, payments=_PAYMENTS, classifications=_CLASSIFICATIONS)
        report = asyncio.run(rs.build_report(_identity(["sell_digital"]), days=7))
        assert "vendas" in report["sections"]
        assert report["sections"]["vendas"]["receita_cents"] == 50000
        assert report["sections"]["vendas"]["pagamentos"] == 2
        assert report["sections"]["vendas"]["fechadas_sem_humano"] == 1
        # Não pediu agendamento nem qualificação
        assert "agenda" not in report["sections"]
        assert "qualificacao" not in report["sections"]

    def test_cliente_que_agenda_ve_agenda(self, monkeypatch):
        _fake_supa(monkeypatch, conversations=_CONVS, classifications=[])
        report = asyncio.run(rs.build_report(_identity(["schedule"]), days=7))
        assert "agenda" in report["sections"]
        assert report["sections"]["agenda"]["agendamentos"] == 1
        assert "vendas" not in report["sections"]

    def test_cliente_que_qualifica_ve_qualificacao(self, monkeypatch):
        _fake_supa(monkeypatch, conversations=_CONVS, classifications=[])
        report = asyncio.run(rs.build_report(_identity(["qualify"]), days=7))
        q = report["sections"]["qualificacao"]
        assert q["leads_com_dados"] == 2   # Maria (dados) + João (nome)
        assert q["enviados_crm"] == 1
        assert q["passados_pro_humano"] == 1
        assert "vendas" not in report["sections"]

    def test_secoes_universais_sempre_presentes(self, monkeypatch):
        _fake_supa(monkeypatch, conversations=_CONVS, classifications=_CLASSIFICATIONS)
        report = asyncio.run(rs.build_report(_identity(["qualify"]), days=7))
        s = report["sections"]
        assert s["atendimento"]["conversas_ativas"] == 3
        assert s["atendimento"]["fora_do_horario"] >= 1  # a das 3h da manhã
        assert s["funil"]["ganhos"] == 1
        assert s["funil"]["perdidos"] == 1
        assert s["follow_up"]["leads_reengajados"] == 2
        assert s["follow_up"]["voltaram_a_negociar"] == 1  # a won com follow-up
        assert s["inteligencia"]["top_assuntos"][0]["tipo"] == "preço"

    def test_formatacao_brl(self):
        assert rs._brl(50000) == "R$ 500,00"
        assert rs._brl(123456) == "R$ 1.234,56"


class TestWhatsAppText:

    def test_texto_por_meta_vendas(self, monkeypatch):
        _fake_supa(monkeypatch, conversations=_CONVS, payments=_PAYMENTS, classifications=_CLASSIFICATIONS)
        identity = _identity(["sell_digital"])
        report = asyncio.run(rs.build_report(identity, days=7))
        msg = rs.format_report_whatsapp(identity, report, "weekly")
        assert "R$ 500,00" in msg
        assert "da semana" in msg
        assert "Clínica Relatório" in msg
        assert "fechadas 100% pela HUMA" in msg
        assert "preço" in msg  # insight

    def test_texto_qualificacao_sem_vendas(self, monkeypatch):
        _fake_supa(monkeypatch, conversations=_CONVS, classifications=[])
        identity = _identity(["qualify"])
        report = asyncio.run(rs.build_report(identity, days=1))
        msg = rs.format_report_whatsapp(identity, report, "daily")
        assert "qualificados" in msg
        assert "R$" not in msg  # não vende → não fala de receita
        assert "de ontem" in msg


class TestOwnerReportJob:

    def _patch_hour(self, monkeypatch, hour):
        class FakeDT(datetime):
            @classmethod
            def utcnow(cls):
                return datetime(2026, 7, 6, hour, 5, 0)

        monkeypatch.setattr(rs, "datetime", FakeDT)

    def test_fora_da_hora_do_cliente_nao_envia(self, monkeypatch):
        # 15h UTC = 12h BRT; cliente default (report_hour=8) não é a hora dele
        self._patch_hour(monkeypatch, 15)
        _fake_supa(monkeypatch, clients=[
            {"client_id": "cli_rep", "report_frequency": "weekly", "owner_phone": "5511999998888"},
        ])
        sent = []

        async def ping(): return True
        async def notify_owner(phone, msg, client_id=""):
            sent.append(phone)

        import huma.services.whatsapp_service as wa_mod
        monkeypatch.setattr(rs.cache, "ping", ping)
        monkeypatch.setattr(wa_mod, "notify_owner", notify_owner)

        asyncio.run(rs.run_owner_reports())
        assert sent == []

    def test_hora_personalizada_do_cliente(self, monkeypatch):
        # Cliente escolheu 12h BRT no drawer → envia às 15h UTC
        self._patch_hour(monkeypatch, 15)
        _fake_supa(monkeypatch, clients=[
            {"client_id": "cli_rep", "report_frequency": "daily",
             "owner_phone": "5511999998888", "report_hour": 12},
        ])
        sent = []
        markers = {}

        async def ping(): return True
        async def get_value(key): return None
        async def set_with_ttl(key, value, ttl=0): markers[key] = value

        import huma.services.db_service as db_mod
        import huma.services.whatsapp_service as wa_mod

        async def get_client(cid):
            return _identity(["schedule"])

        async def notify_owner(phone, msg, client_id=""):
            sent.append(phone)
            return "m1"

        async def fake_build(identity, days=7, date_from="", date_to=""):
            return {"sections": {"atendimento": {"conversas_ativas": 3, "conversas_novas": 1, "fora_do_horario": 0},
                                 "funil": {}, "follow_up": {}, "inteligencia": {"top_assuntos": []}}}

        monkeypatch.setattr(rs.cache, "ping", ping)
        monkeypatch.setattr(rs.cache, "get_value", get_value)
        monkeypatch.setattr(rs.cache, "set_with_ttl", set_with_ttl)
        monkeypatch.setattr(db_mod, "get_client", get_client)
        monkeypatch.setattr(wa_mod, "notify_owner", notify_owner)
        monkeypatch.setattr(rs, "build_report", fake_build)

        asyncio.run(rs.run_owner_reports())
        assert sent == ["5511999998888"]

    def test_destinatarios_extras_recebem_tambem(self, monkeypatch):
        self._patch_hour(monkeypatch, 11)
        _fake_supa(monkeypatch, clients=[
            {"client_id": "cli_rep", "report_frequency": "weekly", "owner_phone": "5511999998888"},
        ])
        sent = []

        async def ping(): return True
        async def get_value(key): return None
        async def set_with_ttl(key, value, ttl=0): pass

        import huma.services.db_service as db_mod
        import huma.services.whatsapp_service as wa_mod

        async def get_client(cid):
            ident = _identity(["schedule"])
            return ident.model_copy(update={"report_recipients": ["5511888887777"]})

        async def notify_owner(phone, msg, client_id=""):
            sent.append(phone)
            return "m1"

        async def fake_build(identity, days=7, date_from="", date_to=""):
            return {"sections": {"atendimento": {"conversas_ativas": 3, "conversas_novas": 1, "fora_do_horario": 0},
                                 "funil": {}, "follow_up": {}, "inteligencia": {"top_assuntos": []}}}

        monkeypatch.setattr(rs.cache, "ping", ping)
        monkeypatch.setattr(rs.cache, "get_value", get_value)
        monkeypatch.setattr(rs.cache, "set_with_ttl", set_with_ttl)
        monkeypatch.setattr(db_mod, "get_client", get_client)
        monkeypatch.setattr(wa_mod, "notify_owner", notify_owner)
        monkeypatch.setattr(rs, "build_report", fake_build)

        asyncio.run(rs.run_owner_reports())
        assert sent == ["5511999998888", "5511888887777"]


    def test_frequencia_respeitada_nao_reenvia(self, monkeypatch):
        self._patch_hour(monkeypatch, 11)
        _fake_supa(monkeypatch, clients=[
            {"client_id": "cli_rep", "report_frequency": "weekly", "owner_phone": "5511999998888"},
        ])
        sent = []

        async def ping(): return True
        async def get_value(key):
            return (datetime(2026, 7, 6) - timedelta(days=2)).isoformat()  # enviado há 2 dias

        async def notify_owner(phone, msg, client_id=""):
            sent.append(phone)

        import huma.services.whatsapp_service as wa_mod
        monkeypatch.setattr(rs.cache, "ping", ping)
        monkeypatch.setattr(rs.cache, "get_value", get_value)
        monkeypatch.setattr(wa_mod, "notify_owner", notify_owner)

        asyncio.run(rs.run_owner_reports())
        assert sent == []  # weekly = 7 dias; só passaram 2

    def test_frequency_off_nao_envia(self, monkeypatch):
        self._patch_hour(monkeypatch, 11)
        _fake_supa(monkeypatch, clients=[
            {"client_id": "cli_rep", "report_frequency": "off", "owner_phone": "5511999998888"},
        ])
        sent = []

        async def ping(): return True
        async def notify_owner(phone, msg, client_id=""):
            sent.append(phone)

        import huma.services.whatsapp_service as wa_mod
        monkeypatch.setattr(rs.cache, "ping", ping)
        monkeypatch.setattr(wa_mod, "notify_owner", notify_owner)

        asyncio.run(rs.run_owner_reports())
        assert sent == []


class TestClientDueNow:
    """_client_due_now: gate puro de hora/dia por cliente."""

    def test_hora_default_8h_brt(self):
        row = {"report_frequency": "weekly", "owner_phone": "x"}
        due, dedup = rs._client_due_now(row, datetime(2026, 7, 6, 11, 5))  # 8h BRT
        assert due and dedup == 7
        due2, _ = rs._client_due_now(row, datetime(2026, 7, 6, 12, 5))
        assert not due2

    def test_dia_da_semana_semanal(self):
        # 2026-07-06 é segunda (weekday 0 em BRT)
        row = {"report_frequency": "weekly", "report_day": "0"}
        due, _ = rs._client_due_now(row, datetime(2026, 7, 6, 11, 5))
        assert due
        row2 = {"report_frequency": "weekly", "report_day": "3"}  # quinta
        due2, _ = rs._client_due_now(row2, datetime(2026, 7, 6, 11, 5))
        assert not due2

    def test_dia_do_mes_mensal_com_dedup_ajustado(self):
        row = {"report_frequency": "monthly", "report_day": "6"}
        due, dedup = rs._client_due_now(row, datetime(2026, 7, 6, 11, 5))
        assert due and dedup == 25  # fevereiro não pode engolir um mês
        row2 = {"report_frequency": "monthly", "report_day": "15"}
        due2, _ = rs._client_due_now(row2, datetime(2026, 7, 6, 11, 5))
        assert not due2

    def test_off_nunca(self):
        due, _ = rs._client_due_now({"report_frequency": "off"}, datetime(2026, 7, 6, 11, 5))
        assert not due


class TestReportsEndpoint:

    def test_sem_auth_401(self):
        from fastapi.testclient import TestClient
        from huma.app import app
        resp = TestClient(app).get("/api/clients/cli_rep/reports")
        assert resp.status_code == 401

    def test_com_sessao_devolve_relatorio(self, monkeypatch):
        import huma.core.auth as auth_mod
        from fastapi.testclient import TestClient
        from huma.app import app

        identity = _identity(["sell_digital"])

        async def get_client(cid):
            return identity if cid == "cli_rep" else None

        async def fake_build(client, days=30, date_from="", date_to=""):
            assert days == 30
            return {"client_id": "cli_rep", "sections": {"atendimento": {"conversas_ativas": 1}}}

        monkeypatch.setattr(auth_mod, "get_client", get_client)
        monkeypatch.setattr(auth_mod, "SESSION_SECRET", "segredo-teste")
        monkeypatch.setattr(rs, "build_report", fake_build)

        cookies = {"huma_session": auth_mod.create_session_token("cli_rep")}
        resp = TestClient(app).get("/api/clients/cli_rep/reports?days=30", cookies=cookies)
        assert resp.status_code == 200
        assert resp.json()["client_id"] == "cli_rep"


class TestReportFrequencyValidation:

    def test_valores_validos(self):
        for v in ("daily", "weekly", "biweekly", "monthly", "off"):
            ident = ClientIdentity(client_id="x", business_name="B", report_frequency=v)
            assert ident.report_frequency == v

    def test_valor_invalido_recusado(self):
        with pytest.raises(Exception):
            ClientIdentity(client_id="x", business_name="B", report_frequency="sempre")


# ================================================================
# EXPORTAÇÃO (.xlsx e .pptx)
# ================================================================

_FAKE_REPORT = {
    "period_days": 30,
    "sections": {
        "atendimento": {"conversas_ativas": 47, "conversas_novas": 12, "fora_do_horario": 9},
        "vendas": {"receita_display": "R$ 3.240,00", "receita_cents": 324000,
                   "pagamentos": 6, "ticket_display": "R$ 540,00", "fechadas_sem_humano": 2},
        "funil": {"descoberta": 20, "negociando": 15, "compromissados": 5, "ganhos": 7, "perdidos": 4},
        "follow_up": {"leads_reengajados": 5, "voltaram_a_negociar": 3},
        "inteligencia": {"top_assuntos": [{"tipo": "preço", "vezes": 14}], "objecoes": 5, "perguntas_preco": 14},
    },
}


class TestExport:

    def test_xlsx_gera_e_reabre_com_dados(self):
        from io import BytesIO
        from openpyxl import load_workbook
        from huma.services import export_service

        data = export_service.report_to_xlsx(_identity(["sell_digital"]), _FAKE_REPORT)
        assert data[:2] == b"PK"  # zip válido (xlsx)

        wb = load_workbook(BytesIO(data))
        ws = wb.active
        textos = [str(c.value) for row in ws.iter_rows() for c in row if c.value is not None]
        blob = " | ".join(textos)
        assert "Clínica Relatório" in blob
        assert "R$ 3.240,00" in blob
        assert "Receita confirmada" in blob
        assert "preço" in blob

    def test_pptx_gera_slides_por_secao(self):
        from io import BytesIO
        from pptx import Presentation
        from huma.services import export_service

        data = export_service.report_to_pptx(_identity(["sell_digital"]), _FAKE_REPORT)
        assert data[:2] == b"PK"

        prs = Presentation(BytesIO(data))
        # capa + atendimento + vendas + funil + follow_up + inteligência + fechamento
        assert len(prs.slides) == 7
        textos = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    textos.append(shape.text_frame.text)
        blob = " | ".join(textos)
        assert "Clínica Relatório" in blob
        assert "R$ 3.240,00" in blob
        assert "Vendas" in blob

    def test_secao_ausente_nao_gera_slide(self):
        from io import BytesIO
        from pptx import Presentation
        from huma.services import export_service

        report = {"period_days": 7, "sections": {
            "atendimento": {"conversas_ativas": 3, "conversas_novas": 1, "fora_do_horario": 0},
            "funil": {"descoberta": 2, "negociando": 1, "compromissados": 0, "ganhos": 0, "perdidos": 0},
            "follow_up": {"leads_reengajados": 0, "voltaram_a_negociar": 0},
            "inteligencia": {"top_assuntos": []},
        }}
        data = export_service.report_to_pptx(_identity(["qualify"]), report)
        prs = Presentation(BytesIO(data))
        blob = " | ".join(
            shape.text_frame.text for slide in prs.slides
            for shape in slide.shapes if shape.has_text_frame
        )
        assert "Vendas" not in blob  # sem seção vendas, sem slide de vendas


class TestExportEndpoint:

    def test_sem_auth_401(self):
        from fastapi.testclient import TestClient
        from huma.app import app
        resp = TestClient(app).get("/api/clients/cli_rep/reports/export?format=xlsx")
        assert resp.status_code == 401

    def test_formato_invalido_400(self, monkeypatch):
        import huma.core.auth as auth_mod
        from fastapi.testclient import TestClient
        from huma.app import app

        async def get_client(cid):
            return _identity(["sell_digital"])

        monkeypatch.setattr(auth_mod, "get_client", get_client)
        monkeypatch.setattr(auth_mod, "SESSION_SECRET", "segredo-teste")
        cookies = {"huma_session": auth_mod.create_session_token("cli_rep")}
        resp = TestClient(app).get("/api/clients/cli_rep/reports/export?format=pdf", cookies=cookies)
        assert resp.status_code == 400

    def test_download_xlsx_com_headers(self, monkeypatch):
        import huma.core.auth as auth_mod
        from fastapi.testclient import TestClient
        from huma.app import app

        identity = _identity(["sell_digital"])

        async def get_client(cid):
            return identity if cid == "cli_rep" else None

        async def fake_build(client, days=30, date_from="", date_to=""):
            return _FAKE_REPORT

        monkeypatch.setattr(auth_mod, "get_client", get_client)
        monkeypatch.setattr(auth_mod, "SESSION_SECRET", "segredo-teste")
        monkeypatch.setattr(rs, "build_report", fake_build)

        cookies = {"huma_session": auth_mod.create_session_token("cli_rep")}
        resp = TestClient(app).get(
            "/api/clients/cli_rep/reports/export?format=xlsx&days=30", cookies=cookies,
        )
        assert resp.status_code == 200
        assert "spreadsheetml" in resp.headers["content-type"]
        assert 'huma-relatorio-30d.xlsx' in resp.headers["content-disposition"]
        assert resp.content[:2] == b"PK"
